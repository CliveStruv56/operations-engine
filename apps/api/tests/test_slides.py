"""Slide export: markdown → Deck parsing, PPTX round-trip, endpoint guards."""

from io import BytesIO
from uuid import uuid4

from pptx import Presentation

from app.db import db
from app.slides import chart_spec, deck_filename, parse_deck, render_pptx
from app.storage import Storage, storage
from tests.conftest import auth, seed_tenant

SAMPLE = """# Forced Quit on Apple Mac M3

## Slide 1 — Why Use Forced Quit?
- Application becomes unresponsive or freezes
- **Preserves** other open programs [1]

**Speaker notes:** Emphasize that forced quit is a last-resort tool;
encourage normal quit first.

---

## Slide 2 — Accessing the Force Quit Menu
- Click the **Apple menu** → **Force Quit…**
- Or press Option-Command-Escape
"""


def test_parse_deck_structure():
    deck = parse_deck(SAMPLE)
    assert deck is not None
    assert deck.title == "Forced Quit on Apple Mac M3"
    assert [s.title for s in deck.slides] == [
        "Why Use Forced Quit?",
        "Accessing the Force Quit Menu",
    ]
    # Markdown markup and citation markers are stripped from bullets.
    assert deck.slides[0].bullets[1] == "Preserves other open programs"
    # Multi-line notes paragraphs are joined.
    assert deck.slides[0].notes is not None
    assert deck.slides[0].notes.startswith("Emphasize")
    assert deck.slides[0].notes.endswith("normal quit first.")
    assert deck.slides[1].notes is None


def test_parse_rejects_prose():
    assert parse_deck("Just an ordinary chat answer with no headings.") is None
    assert parse_deck("# A title but no slides\n\nSome prose.") is None


def test_parse_deck_title_falls_back_to_first_slide():
    deck = parse_deck("## Slide 1 — Only Slide\n- point")
    assert deck is not None
    assert deck.title == "Only Slide"


def test_render_pptx_roundtrip():
    deck = parse_deck(SAMPLE)
    data = render_pptx(deck, "#336699")
    assert data[:2] == b"PK"  # zip container
    prs = Presentation(BytesIO(data))
    assert len(prs.slides) == 3  # title + 2 content
    notes = prs.slides[1].notes_slide.notes_text_frame.text
    assert "last-resort" in notes


def test_chart_spec_detection():
    spec = chart_spec(["Q1: £1.2m", "Q2: £3.4m", "Q3 — £2.0m"])
    assert spec is not None
    assert spec.categories == ["Q1", "Q2", "Q3"]
    assert spec.values == [1.2, 3.4, 2.0]
    assert spec.unit == "£m"

    # Trailing unit words count too.
    spec = chart_spec(["Phase 1: 40 homes", "Phase 2: 60 homes", "Phase 3: 50 homes"])
    assert spec is not None
    assert spec.unit == "homes"

    # Below the row threshold, mixed prose, or year-only timelines → bullets.
    assert chart_spec(["Q1: 10", "Q2: 20"]) is None
    assert chart_spec(["Q1: 10", "Q2: 20", "Click the Apple menu"]) is None
    assert chart_spec(["Start: 2024", "Review: 2025", "Finish: 2026"]) is None


CHART_DECK = """# Delivery numbers

## Slide 1 — Homes per phase
- Phase 1: 40 homes
- Phase 2: 60 homes
- Phase 3: 50 homes
"""


def test_render_chart_slide():
    deck = parse_deck(CHART_DECK)
    data = render_pptx(deck)
    prs = Presentation(BytesIO(data))
    charts = [s for s in prs.slides[1].shapes if s.has_chart]
    assert len(charts) == 1
    chart = charts[0].chart
    plot = chart.plots[0]
    assert list(plot.categories) == ["Phase 1", "Phase 2", "Phase 3"]
    assert [v for v in plot.series[0].values] == [40.0, 60.0, 50.0]


def _make_template(sample_slides: int = 2) -> bytes:
    """A stand-in corporate master: python-pptx's default template has a
    Title Slide layout and a Title-and-Content layout, plus sample slides
    that the renderer must clear out."""
    prs = Presentation()
    for _ in range(sample_slides):
        prs.slides.add_slide(prs.slide_layouts[0])
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def test_render_on_template_uses_master_layouts():
    deck = parse_deck(SAMPLE)
    data = render_pptx(deck, template=_make_template())
    prs = Presentation(BytesIO(data))
    # Sample slides cleared: title + 2 content only.
    assert len(prs.slides) == 3
    # Placeholder-based rendering: title placeholders exist and carry text
    # (the generated theme uses raw text boxes, so .title would be None).
    assert prs.slides[0].shapes.title.text == "Forced Quit on Apple Mac M3"
    assert prs.slides[1].shapes.title.text == "Why Use Forced Quit?"
    assert "last-resort" in prs.slides[1].notes_slide.notes_text_frame.text


def test_render_falls_back_on_broken_template():
    deck = parse_deck(SAMPLE)
    data = render_pptx(deck, template=b"not a pptx at all")
    prs = Presentation(BytesIO(data))
    assert len(prs.slides) == 3  # generated theme still delivered


def test_render_chart_on_template():
    deck = parse_deck(CHART_DECK)
    data = render_pptx(deck, template=_make_template())
    prs = Presentation(BytesIO(data))
    charts = [s for s in prs.slides[1].shapes if s.has_chart]
    assert len(charts) == 1


def test_deck_filename():
    assert deck_filename("Forced Quit on Apple Mac M3!") == "forced-quit-on-apple-mac-m3.pptx"
    assert deck_filename("***") == "slides.pptx"


async def _seed_deck_message(tenant, content: str = SAMPLE) -> str:
    async with db.tenant_tx(tenant.owner_id, tenant.id) as conn:
        return str(
            await conn.fetchval(
                "insert into messages (tenant_id, conversation_id, role, content)"
                " values ($1, $2, 'assistant', $3) returning id",
                tenant.id,
                tenant.conversation_id,
                content,
            )
        )


async def test_export_requires_storage(client):
    tenant = await seed_tenant(client, f"pptx503-{uuid4().hex[:6]}")
    message_id = await _seed_deck_message(tenant)
    resp = await client.post(
        f"/api/v1/conversations/{tenant.conversation_id}/messages/{message_id}/slides",
        headers=auth(tenant.owner_id, tenant.id),
    )
    assert resp.status_code == 503
    assert resp.json()["error"]["code"] == "storage_unavailable"


async def test_export_flow_and_guards(client, monkeypatch):
    tenant = await seed_tenant(client, f"pptx-{uuid4().hex[:6]}")
    headers = auth(tenant.owner_id, tenant.id)
    deck_msg = await _seed_deck_message(tenant)
    prose_msg = await _seed_deck_message(tenant, "Just a normal answer.")

    uploads: dict[str, tuple[bytes, str]] = {}

    async def fake_upload(key, data, mime):
        uploads[key] = (data, mime)

    monkeypatch.setattr(Storage, "enabled", property(lambda self: True))
    monkeypatch.setattr(storage, "upload_bytes", fake_upload)
    monkeypatch.setattr(storage, "presign_get", lambda key: f"https://signed.example/{key}")

    resp = await client.post(
        f"/api/v1/conversations/{tenant.conversation_id}/messages/{deck_msg}/slides",
        headers=headers,
    )
    assert resp.status_code == 200, resp.text
    body = resp.json()
    assert body["slide_count"] == 2
    assert body["filename"] == "forced-quit-on-apple-mac-m3.pptx"
    assert body["download_url"] == f"https://signed.example/{tenant.id}/slides/{deck_msg}.pptx"
    data, mime = uploads[f"{tenant.id}/slides/{deck_msg}.pptx"]
    assert data[:2] == b"PK"
    assert mime.endswith("presentationml.presentation")

    # Non-deck message → 400.
    resp = await client.post(
        f"/api/v1/conversations/{tenant.conversation_id}/messages/{prose_msg}/slides",
        headers=headers,
    )
    assert resp.status_code == 400
    assert resp.json()["error"]["code"] == "not_a_deck"

    # Cross-tenant direct-object reference → 404 under the attacker's context.
    other = await seed_tenant(client, f"pptxb-{uuid4().hex[:6]}")
    resp = await client.post(
        f"/api/v1/conversations/{tenant.conversation_id}/messages/{deck_msg}/slides",
        headers=auth(other.owner_id, other.id),
    )
    assert resp.status_code == 404

    # Audit written.
    async with db.tenant_tx(tenant.owner_id, tenant.id) as conn:
        actions = {
            r["action"]
            for r in await conn.fetch(
                "select action from audit_log where tenant_id = $1", tenant.id
            )
        }
    assert "message.slides_export" in actions


async def test_export_uses_tenant_template(client, monkeypatch):
    tenant = await seed_tenant(client, f"pptxt-{uuid4().hex[:6]}")
    headers = auth(tenant.owner_id, tenant.id)
    deck_msg = await _seed_deck_message(tenant)
    template = _make_template()

    uploads: dict[str, bytes] = {}

    async def fake_upload(key, data, mime):
        uploads[key] = data

    async def fake_download(key):
        assert key == f"{tenant.id}/brand/slides-template.pptx"
        return template

    monkeypatch.setattr(Storage, "enabled", property(lambda self: True))
    monkeypatch.setattr(storage, "upload_bytes", fake_upload)
    monkeypatch.setattr(storage, "download_bytes", fake_download)
    monkeypatch.setattr(storage, "presign_get", lambda key: f"https://signed.example/{key}")

    resp = await client.patch(
        "/api/v1/tenants/me",
        json={"brand": {"slides_template_key": f"{tenant.id}/brand/slides-template.pptx"}},
        headers=headers,
    )
    assert resp.status_code == 200

    resp = await client.post(
        f"/api/v1/conversations/{tenant.conversation_id}/messages/{deck_msg}/slides",
        headers=headers,
    )
    assert resp.status_code == 200, resp.text
    prs = Presentation(BytesIO(uploads[f"{tenant.id}/slides/{deck_msg}.pptx"]))
    # Placeholder-based title proves the master's layouts were used.
    assert prs.slides[0].shapes.title.text == "Forced Quit on Apple Mac M3"
