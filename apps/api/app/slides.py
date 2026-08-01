"""Slide-deck export: parse a slides-mode chat message (structured markdown,
shape enforced by prompts.TASK_PROMPTS["slides"]) and render it as a native,
editable PPTX. Parsing and rendering are deterministic and sub-second — no
LLM call, no job queue.

Theming, in order of preference:
1. Tenant slides template (brand.slides_template_key — a corporate .pptx
   master): slides are built on its layouts/placeholders so fonts, colours
   and list styles all come from the master.
2. Generated theme: blank 16:9 deck styled with the tenant accent + logo.

Slides whose bullets are all "Label: number" pairs render as a native
column chart instead of a list.

python-pptx is MIT-licensed (constraint 1).
"""

import re
from dataclasses import dataclass
from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

DEFAULT_ACCENT = "#1f6d53"
INK = RGBColor(0x1C, 0x1E, 0x21)

_HEX = re.compile(r"^#[0-9a-fA-F]{6}$")
_H1 = re.compile(r"^#\s+(.+)$")
_H2 = re.compile(r"^##\s+(.+)$")
_BULLET = re.compile(r"^\s*[-*]\s+(.+)$")
_NOTES = re.compile(r"^\**Speaker notes:?\**\s*(.*)$", re.IGNORECASE)
_SLIDE_PREFIX = re.compile(r"^Slide\s*\d+\s*[—–-]\s*", re.IGNORECASE)
# Inline markdown + resolved citation markers ([1]) don't belong on slides.
_MARKUP = re.compile(r"\*\*|\*|`|\[\d+\]")

# "Label: 42", "Q1 — £3.2m", "Homes delivered = 150 units"
_DATA_BULLET = re.compile(
    r"^(?P<label>.+?)\s*[:=–—-]\s*(?P<cur>[£$€])?\s*"
    r"(?P<num>\d[\d,]*(?:\.\d+)?)\s*(?P<suffix>%|bn|m|k)?\s*(?P<rest>[A-Za-z][A-Za-z ]{0,18})?$"
)
MIN_CHART_ROWS = 3


@dataclass(frozen=True)
class Slide:
    title: str
    bullets: list[str]
    notes: str | None


@dataclass(frozen=True)
class Deck:
    title: str
    slides: list[Slide]


@dataclass(frozen=True)
class ChartSpec:
    categories: list[str]
    values: list[float]
    unit: str


def _clean(text: str) -> str:
    return re.sub(r"\s{2,}", " ", _MARKUP.sub("", text)).strip()


def parse_deck(content: str) -> Deck | None:
    """Deck from slides-mode markdown, or None when the message isn't one.
    Tolerant of drift: any ## heading starts a slide; bullets are - / * lines;
    a "Speaker notes:" paragraph becomes the notes pane."""
    deck_title: str | None = None
    slides: list[Slide] = []
    title: str | None = None
    bullets: list[str] = []
    notes_parts: list[str] = []
    in_notes = False

    def _flush() -> None:
        nonlocal title, bullets, notes_parts, in_notes
        if title is not None:
            slides.append(Slide(title, bullets, " ".join(notes_parts) or None))
        title, bullets, notes_parts, in_notes = None, [], [], False

    for raw in content.splitlines():
        line = raw.strip()
        if not line or line == "---":
            in_notes = False
            continue
        if (m := _H2.match(line)) is not None:
            _flush()
            title = _clean(_SLIDE_PREFIX.sub("", m.group(1)))
            continue
        if (m := _H1.match(line)) is not None:
            if deck_title is None:
                deck_title = _clean(m.group(1))
            continue
        if title is None:
            continue  # prose before the first slide
        if (m := _NOTES.match(line)) is not None:
            in_notes = True
            if m.group(1):
                notes_parts.append(_clean(m.group(1)))
            continue
        if (m := _BULLET.match(raw)) is not None:
            in_notes = False
            bullets.append(_clean(m.group(1)))
            continue
        if in_notes:
            notes_parts.append(_clean(line))

    _flush()
    if not slides:
        return None
    return Deck(deck_title or slides[0].title, slides)


def chart_spec(bullets: list[str]) -> ChartSpec | None:
    """A slide charts instead of listing when every bullet is a label→number
    pair (≥ MIN_CHART_ROWS). All-year values with no unit stay as bullets —
    "Start: 2024 / Finish: 2026" is a timeline, not a series."""
    if len(bullets) < MIN_CHART_ROWS:
        return None
    categories: list[str] = []
    values: list[float] = []
    units: list[str] = []
    for bullet in bullets:
        m = _DATA_BULLET.match(bullet)
        if m is None:
            return None
        categories.append(m.group("label").strip())
        values.append(float(m.group("num").replace(",", "")))
        unit = "".join(filter(None, (m.group("cur"), m.group("suffix")))) or (
            (m.group("rest") or "").strip()
        )
        if unit:
            units.append(unit)
    if not units and all(v.is_integer() and 1900 <= v <= 2100 for v in values):
        return None
    unit = max(set(units), key=units.count) if units else "Value"
    return ChartSpec(categories, values, unit)


# --- Rendering ---------------------------------------------------------------


def _accent_color(accent_hex: str | None) -> RGBColor:
    hex_value = accent_hex if accent_hex and _HEX.match(accent_hex) else DEFAULT_ACCENT
    return RGBColor.from_string(hex_value.lstrip("#"))


def _add_rule(slide, accent: RGBColor, left, top, width, height=Inches(0.05)) -> None:
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent
    rule.line.fill.background()
    rule.shadow.inherit = False


def _add_chart(slide, spec: ChartSpec, rect, accent: RGBColor | None) -> None:
    """Native column chart. accent=None (template mode) leaves the series on
    the master's theme colours."""
    data = CategoryChartData()
    data.categories = spec.categories
    data.add_series(spec.unit, spec.values)
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *rect, data)
    chart = frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    if accent is not None:
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = accent


def _set_notes(slide, notes: str | None) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# --- Template mode -----------------------------------------------------------

_TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
_BODY_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)


def _clear_slides(prs) -> None:
    """Corporate masters arrive as .pptx files full of sample slides; only
    our generated slides should survive. (python-pptx has no public delete.)"""
    id_list = prs.slides._sldIdLst  # noqa: SLF001 — the established recipe
    for slide_id in list(id_list):
        prs.part.drop_rel(slide_id.get(qn("r:id")))
        id_list.remove(slide_id)


def _placeholder(slide, types):
    for ph in slide.placeholders:
        if ph.placeholder_format.type in types:
            return ph
    return None


def _pick_layouts(prs) -> tuple | None:
    """(title_layout, content_layout) from the master, or None when the
    template has no usable placeholder layouts."""
    title_layout = None
    content_layout = None
    for layout in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in layout.placeholders}
        has_title = any(t in types for t in _TITLE_TYPES)
        has_body = any(t in types for t in _BODY_TYPES)
        if title_layout is None and PP_PLACEHOLDER.CENTER_TITLE in types:
            title_layout = layout
        if content_layout is None and has_title and has_body:
            content_layout = layout
    if title_layout is None:
        title_layout = content_layout
    if title_layout is None or content_layout is None:
        return None
    return title_layout, content_layout


def _render_on_template(deck: Deck, template: bytes) -> bytes | None:
    """Build the deck on the tenant's own master; None → caller falls back to
    the generated theme (unreadable file or no usable layouts)."""
    try:
        prs = Presentation(BytesIO(template))
        layouts = _pick_layouts(prs)
        if layouts is None:
            return None
        title_layout, content_layout = layouts
        _clear_slides(prs)

        slide = prs.slides.add_slide(title_layout)
        title_ph = _placeholder(slide, _TITLE_TYPES)
        if title_ph is not None:
            title_ph.text = deck.title

        for item in deck.slides:
            slide = prs.slides.add_slide(content_layout)
            title_ph = _placeholder(slide, _TITLE_TYPES)
            if title_ph is not None:
                title_ph.text = item.title
            body = _placeholder(slide, _BODY_TYPES)
            spec = chart_spec(item.bullets)
            if spec is not None and body is not None:
                rect = (body.left, body.top, body.width, body.height)
                body._element.getparent().remove(body._element)  # noqa: SLF001
                _add_chart(slide, spec, rect, accent=None)  # master theme colours
            elif body is not None and item.bullets:
                frame = body.text_frame
                for i, bullet in enumerate(item.bullets):
                    para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                    para.text = bullet  # master list styles supply the bullets
            _set_notes(slide, item.notes)

        out = BytesIO()
        prs.save(out)
        return out.getvalue()
    except Exception:
        # A broken upload must not block the export — presigned PUTs mean the
        # API never saw the template bytes at upload time to validate them.
        return None


# --- Generated theme ---------------------------------------------------------


def _render_generated(deck: Deck, accent_hex: str | None, logo: bytes | None) -> bytes:
    accent = _accent_color(accent_hex)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    if logo:
        slide.shapes.add_picture(BytesIO(logo), Inches(0.6), Inches(0.5), height=Inches(0.5))
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8))
    box.text_frame.word_wrap = True
    run = box.text_frame.paragraphs[0].add_run()
    run.text = deck.title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = INK
    _add_rule(slide, accent, Inches(0.95), Inches(4.35), Inches(2.5), Inches(0.07))

    for item in deck.slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.9))
        box.text_frame.word_wrap = True
        run = box.text_frame.paragraphs[0].add_run()
        run.text = item.title
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = accent
        _add_rule(slide, accent, Inches(0.75), Inches(1.25), Inches(1.6))

        spec = chart_spec(item.bullets)
        if spec is not None:
            rect = (Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.2))
            _add_chart(slide, spec, rect, accent)
        elif item.bullets:
            box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.2))
            frame = box.text_frame
            frame.word_wrap = True
            for i, bullet in enumerate(item.bullets):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                para.space_after = Pt(10)
                run = para.add_run()
                run.text = f"•  {bullet}"
                run.font.size = Pt(17)
                run.font.color.rgb = INK
        _set_notes(slide, item.notes)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def render_pptx(
    deck: Deck,
    accent_hex: str | None = None,
    logo: bytes | None = None,
    template: bytes | None = None,
) -> bytes:
    if template is not None:
        rendered = _render_on_template(deck, template)
        if rendered is not None:
            return rendered
    return _render_generated(deck, accent_hex, logo)


def deck_filename(title: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", title).strip("-").lower()[:60]
    return f"{slug or 'slides'}.pptx"
