#!/usr/bin/env python3
"""Generate Flowgrid OS pitch decks from the one-pager content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent

# Hearth palette
TERRACOTTA = RGBColor(0xB1, 0x4E, 0x2E)
TERRACOTTA_DARK = RGBColor(0x8A, 0x3A, 0x22)
CREAM = RGBColor(0xF6, 0xF2, 0xED)
INK = RGBColor(0x2D, 0x2A, 0x26)
INK_SOFT = RGBColor(0x5B, 0x57, 0x52)
SAGE = RGBColor(0x6B, 0x8E, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DISPLAY_FONT = "Georgia"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def _set_text(run, text, size, bold=False, color=INK, font=BODY_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT, font=BODY_FONT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    _set_text(p.runs[0], text, size, bold=bold, color=color, font=font)
    return shape


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, line_spacing=1.2):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = BODY_FONT
        p.space_after = Pt(10)
        p.line_spacing = line_spacing
    return shape


def add_two_col_bullets(slide, left, top, width, height, col1_items, col2_items, size=15):
    col_w = (width - Inches(0.4)) / 2
    add_bullets(slide, left, top, col_w, height, col1_items, size=size)
    add_bullets(slide, left + col_w + Inches(0.4), top, col_w, height, col2_items, size=size)


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()
    add_textbox(slide, MARGIN, Inches(0.22), SLIDE_W - MARGIN * 2, Inches(0.8), title, size=34, bold=True, color=WHITE, font=DISPLAY_FONT)


def add_footer(slide, text="flowgridos.co.uk"):
    add_textbox(slide, MARGIN, SLIDE_H - Inches(0.35), SLIDE_W - MARGIN * 2, Inches(0.25), text, size=10, color=INK_SOFT, align=PP_ALIGN.RIGHT)


def title_slide(prs, headline, subhead, footer_text="flowgridos.co.uk"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()

    # brand mark
    mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.6), Inches(0.55), Inches(0.55))
    mark.fill.solid()
    mark.fill.fore_color.rgb = TERRACOTTA
    mark.line.fill.background()

    add_textbox(slide, MARGIN + Inches(0.8), Inches(1.55), Inches(10), Inches(0.7), "Flowgrid OS", size=26, bold=True, color=INK, font=DISPLAY_FONT)

    add_textbox(slide, MARGIN, Inches(2.5), Inches(11.5), Inches(1.4), headline, size=48, bold=True, color=TERRACOTTA_DARK, font=DISPLAY_FONT)
    add_textbox(slide, MARGIN, Inches(4.0), Inches(11.5), Inches(1.0), subhead, size=24, color=INK_SOFT)

    add_footer(slide, footer_text)
    return slide


def section_slide(prs, title, bullets, footer_text="flowgridos.co.uk"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    add_bullets(slide, MARGIN, Inches(1.45), SLIDE_W - MARGIN * 2, Inches(5.6), bullets, size=22)
    add_footer(slide, footer_text)
    return slide


def customer_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "AI workspace for UK small businesses", "One workspace. Your brand. Your data.")

    section_slide(prs, "The problem", [
        "SMBs are buying AI in fragments — chat, documents, slides, research, each with its own login and bill.",
        "Closed models are expensive and often train on your data.",
        "Company knowledge sits in folders and spreadsheets, never where the work happens.",
    ])

    section_slide(prs, "The solution", [
        "A single branded workspace for every client: their logo, their accent, their documents.",
        "Answers cite the source page, so staff can trust what the AI says.",
        "Every task is routed to the cheapest capable open-weight model — no closed-model lock-in.",
        "Usage caps and per-tenant budgets prevent surprise invoices.",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What you get")
    add_two_col_bullets(
        slide, MARGIN, Inches(1.5), SLIDE_W - MARGIN * 2, Inches(5.4),
        [
            "Branded AI chat with page-level citations",
            "Cited knowledge vault (PDF, Word, Excel, Markdown)",
            "Research mode with cited web sources",
            "Native PowerPoint slide export",
        ],
        [
            "CRM contact book with CSV import",
            "Project spine for stage-gated work",
            "Email & calendar sync (roadmap)",
            "Meeting intelligence (roadmap)",
        ],
        size=19,
    )
    add_footer(slide)

    section_slide(prs, "Privacy & trust", [
        "Tenant isolation enforced by Postgres row-level security.",
        "Zero-data-retention US/EU hosts by default; nothing trains on your files.",
        "Soft caps nudge, hard caps protect; full audit logs.",
        "MIT/Apache stack only — no hidden licensing overhang.",
    ])

    # Pricing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Simple, per-seat pricing")
    box_w = (SLIDE_W - MARGIN * 2 - Inches(0.6)) / 3
    boxes = [
        ("Core", "£29", "/seat/mo · min 3", [
            "Workspace & chat",
            "Vault & citations",
            "Research & slides",
            "CRM contacts",
            "Standard support",
        ]),
        ("Pro", "£49", "/seat/mo", [
            "Everything in Core",
            "Premium reasoning",
            "Deep research",
            "Advanced reporting",
            "Priority support",
        ]),
        ("Managed", "£799", "/mo + £1,500 onboarding", [
            "Includes 10 Pro seats",
            "White-glove onboarding",
            "Custom workflows",
            "Quarterly reviews",
            "+£39/extra seat",
        ]),
    ]
    for i, (name, price, suffix, items) in enumerate(boxes):
        left = MARGIN + i * (box_w + Inches(0.3))
        top = Inches(1.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CREAM if name != "Pro" else TERRACOTTA
        shape.line.color.rgb = TERRACOTTA
        shape.line.width = Pt(2)

        fill = WHITE if name == "Pro" else INK
        accent = WHITE if name == "Pro" else TERRACOTTA

        add_textbox(slide, left + Inches(0.15), top + Inches(0.15), box_w - Inches(0.3), Inches(0.5), name, size=22, bold=True, color=accent, align=PP_ALIGN.CENTER, font=DISPLAY_FONT)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.7), box_w - Inches(0.3), Inches(0.6), price, size=32, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), top + Inches(1.25), box_w - Inches(0.3), Inches(0.4), suffix, size=12, color=fill, align=PP_ALIGN.CENTER)
        add_bullets(slide, left + Inches(0.2), top + Inches(1.75), box_w - Inches(0.4), Inches(2.8), items, size=14, color=fill)
    add_footer(slide)

    # CTA slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TERRACOTTA
    bg.line.fill.background()
    add_textbox(slide, MARGIN, Inches(2.2), SLIDE_W - MARGIN * 2, Inches(1.2), "See your workspace before you commit.", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=DISPLAY_FONT)
    add_textbox(slide, MARGIN, Inches(3.5), SLIDE_W - MARGIN * 2, Inches(0.8), "Book a 20-minute demo and we’ll build a live workspace around your own documents.", size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Inches(4.5), SLIDE_W - MARGIN * 2, Inches(0.6), "hello@flowgridos.co.uk · 14-day trial, no card required", size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, "flowgridos.co.uk")

    path = ROOT / "flowgrid-os-customer-deck.pptx"
    prs.save(path)
    print(f"Saved {path}")


def partner_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Partner programme", "Sell a branded AI workspace without building one.", footer_text="flowgridos.co.uk/partners")

    section_slide(prs, "The opportunity", [
        "SMBs want one AI bill and their own brand — not a patchwork of tools.",
        "Agencies, MSPs and consultancies can own the client relationship.",
        "Flowgrid OS handles infrastructure, model routing, security and updates.",
    ], footer_text="flowgridos.co.uk/partners")

    section_slide(prs, "The platform you resell", [
        "Brand layer — tenant logo, accent and white-label shell.",
        "Model router — cost-routed open weights with automatic fallbacks.",
        "Vault & RAG — cited retrieval with per-tenant database isolation.",
        "Modules — chat, research, slides, CRM and project spine.",
    ], footer_text="flowgridos.co.uk/partners")

    section_slide(prs, "Why partner", [
        "High recurring margins: ~89% Core, ~88% Pro, ~70% Managed after support labour.",
        "Zero licensing overhang: MIT/Apache stack end-to-end.",
        "Sticky clients: vault, projects and usage data increase switching costs.",
        "Fast to deploy: new branded workspaces provisioned in minutes.",
    ], footer_text="flowgridos.co.uk/partners")

    # Economics / pricing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Partner economics")
    box_w = (SLIDE_W - MARGIN * 2 - Inches(0.6)) / 3
    boxes = [
        ("Core", "£29", "/seat/mo · min 3", "~89% gross margin", [
            "Workspace, vault & chat",
            "Research & slide export",
            "CRM contacts",
            "Self-serve support",
        ]),
        ("Pro", "£49", "/seat/mo", "~88% gross margin", [
            "Everything in Core",
            "Premium reasoning",
            "Deep research",
            "Advanced reporting",
        ]),
        ("Managed", "£799", "/mo + £1,500 onboarding", "~70% margin after labour", [
            "Includes 10 Pro seats",
            "White-glove onboarding",
            "Quarterly reviews & SLA",
            "+£39/extra seat",
        ]),
    ]
    for i, (name, price, suffix, margin_text, items) in enumerate(boxes):
        left = MARGIN + i * (box_w + Inches(0.3))
        top = Inches(1.45)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CREAM if name != "Pro" else TERRACOTTA
        shape.line.color.rgb = TERRACOTTA
        shape.line.width = Pt(2)

        fill = WHITE if name == "Pro" else INK
        accent = WHITE if name == "Pro" else TERRACOTTA

        add_textbox(slide, left + Inches(0.15), top + Inches(0.15), box_w - Inches(0.3), Inches(0.5), name, size=22, bold=True, color=accent, align=PP_ALIGN.CENTER, font=DISPLAY_FONT)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.7), box_w - Inches(0.3), Inches(0.6), price, size=32, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), top + Inches(1.25), box_w - Inches(0.3), Inches(0.35), suffix, size=12, color=fill, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), top + Inches(1.6), box_w - Inches(0.3), Inches(0.4), margin_text, size=15, bold=True, color=SAGE if name != "Pro" else WHITE, align=PP_ALIGN.CENTER)
        add_bullets(slide, left + Inches(0.2), top + Inches(2.15), box_w - Inches(0.4), Inches(2.7), items, size=14, color=fill)
    add_footer(slide, "flowgridos.co.uk/partners")

    # How it works
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "How it works")
    steps = [
        ("1", "Sign a client", "Pick tier & collect brand assets."),
        ("2", "Provision workspace", "Operator console creates tenant, keys & invites."),
        ("3", "Curate & train", "Upload docs, run onboarding, set caps."),
        ("4", "Bill monthly", "You invoice client; we invoice you for COGS."),
    ]
    box_w = (SLIDE_W - MARGIN * 2 - Inches(0.9)) / 4
    for i, (num, title, desc) in enumerate(steps):
        left = MARGIN + i * (box_w + Inches(0.3))
        top = Inches(1.9)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, Inches(3.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CREAM
        shape.line.color.rgb = TERRACOTTA
        shape.line.width = Pt(2)

        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + box_w / 2 - Inches(0.3), top + Inches(0.25), Inches(0.6), Inches(0.6))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = TERRACOTTA
        num_shape.line.fill.background()
        add_textbox(slide, left + box_w / 2 - Inches(0.3), top + Inches(0.32), Inches(0.6), Inches(0.5), num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, left + Inches(0.15), top + Inches(1.1), box_w - Inches(0.3), Inches(0.5), title, size=17, bold=True, color=TERRACOTTA_DARK, align=PP_ALIGN.CENTER, font=DISPLAY_FONT)
        add_textbox(slide, left + Inches(0.15), top + Inches(1.7), box_w - Inches(0.3), Inches(1.4), desc, size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_footer(slide, "flowgridos.co.uk/partners")

    # CTA slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TERRACOTTA
    bg.line.fill.background()
    add_textbox(slide, MARGIN, Inches(2.2), SLIDE_W - MARGIN * 2, Inches(1.2), "Become a Flowgrid OS partner.", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=DISPLAY_FONT)
    add_textbox(slide, MARGIN, Inches(3.5), SLIDE_W - MARGIN * 2, Inches(0.8), "Book a partner call and we’ll share the operator console, margin model and onboarding runbook.", size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Inches(4.5), SLIDE_W - MARGIN * 2, Inches(0.6), "partners@flowgridos.co.uk", size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, "flowgridos.co.uk/partners")

    path = ROOT / "flowgrid-os-partner-deck.pptx"
    prs.save(path)
    print(f"Saved {path}")


if __name__ == "__main__":
    customer_deck()
    partner_deck()
